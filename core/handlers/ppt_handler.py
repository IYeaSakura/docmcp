"""
PowerPoint文档处理器模块

支持 .ppt 和 .pptx 格式的PowerPoint文档处理。
"""

import io
import logging
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, BinaryIO, Dict, Iterator, List, Optional, Tuple, Union

from .base import BaseDocument, BaseDocumentHandler, BatchProcessor
from ..document import DocumentMetadata, DocumentType, ExtractedContent
from ..utils import TempFileManager, detect_file_type

# 配置日志
logger = logging.getLogger(__name__)

# 尝试导入依赖库
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx 未安装，PPTX功能受限")


try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False
    logger.warning("Pillow 未安装，图片处理功能受限")


class PowerPointDocument(BaseDocument):
    """
    PowerPoint文档对象
    
    支持 .ppt 和 .pptx 格式的PowerPoint文档。
    """
    
    def __init__(
        self,
        file_path: Optional[Union[str, Path]] = None,
        file_stream: Optional[BinaryIO] = None,
        document_type: Optional[DocumentType] = None
    ):
        """
        初始化PowerPoint文档对象
        
        Args:
            file_path: 文档文件路径
            file_stream: 文档文件流
            document_type: 文档类型
        """
        super().__init__(file_path, file_stream, document_type)
        self._presentation: Optional[Any] = None
        self._temp_manager = TempFileManager()
        self._slide_count: int = 0
    
    def load(self, **kwargs) -> "PowerPointDocument":
        """
        加载PowerPoint文档
        
        Returns:
            self: 支持链式调用
        """
        if self._is_loaded:
            return self
        
        try:
            # 确定文档类型
            if self._document_type == DocumentType.UNKNOWN:
                if self._file_path:
                    self._document_type = detect_file_type(self._file_path)
                elif self._file_stream:
                    self._document_type = detect_file_type(self._file_stream)
            
            # 加载文档
            if self._document_type == DocumentType.PPTX:
                self._load_pptx()
            elif self._document_type == DocumentType.PPT:
                self._load_ppt()
            else:
                self._try_auto_load()
            
            self._is_loaded = True
            self._update_metadata_from_file()
            
        except Exception as e:
            logger.error(f"加载PowerPoint文档失败: {e}")
            raise ValueError(f"无法加载PowerPoint文档: {e}")
        
        return self
    
    def _load_pptx(self) -> None:
        """加载PPTX格式文档"""
        if not PPTX_AVAILABLE:
            raise ImportError("需要安装 python-pptx 库: pip install python-pptx")
        
        if self._file_stream:
            self._file_stream.seek(0)
            self._presentation = Presentation(self._file_stream)
        elif self._file_path:
            self._presentation = Presentation(self._file_path)
        else:
            raise ValueError("没有可用的文件源")
        
        self._slide_count = len(self._presentation.slides)
    
    def _load_ppt(self) -> None:
        """加载PPT格式文档（旧版）"""
        # 旧版PPT需要转换为PPTX
        try:
            pptx_path = self._convert_ppt_to_pptx()
            self._presentation = Presentation(pptx_path)
            self._slide_count = len(self._presentation.slides)
            self._document_type = DocumentType.PPTX  # 转换后类型改变
        except Exception as e:
            logger.error(f"加载PPT文件失败: {e}")
            raise
    
    def _convert_ppt_to_pptx(self) -> Path:
        """
        将PPT转换为PPTX
        
        Returns:
            转换后的PPTX文件路径
        """
        if not self._file_path:
            raise ValueError("需要文件路径才能转换PPT")
        
        # 尝试使用LibreOffice转换
        try:
            return self._convert_with_libreoffice()
        except Exception as e:
            logger.warning(f"LibreOffice转换失败: {e}")
        
        raise ValueError("无法转换PPT文件，请安装LibreOffice")
    
    def _convert_with_libreoffice(self) -> Path:
        """使用LibreOffice转换"""
        if not self._file_path:
            raise ValueError("需要文件路径")
        
        output_dir = self._temp_manager.create_temp_dir()
        
        cmd = [
            'soffice',
            '--headless',
            '--convert-to', 'pptx',
            '--outdir', str(output_dir),
            str(self._file_path)
        ]
        
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=120
        )
        
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice转换失败: {result.stderr}")
        
        # 找到转换后的文件
        output_file = output_dir / f"{self._file_path.stem}.pptx"
        if output_file.exists():
            return output_file
        
        # 尝试其他文件名
        for f in output_dir.glob("*.pptx"):
            return f
        
        raise FileNotFoundError("转换后的文件未找到")
    
    def get_slide_count(self) -> int:
        """
        获取幻灯片数量
        
        Returns:
            幻灯片数量
        """
        if not self._is_loaded:
            self.load()
        return self._slide_count
    
    def extract_text(self, **kwargs) -> str:
        """
        提取纯文本内容
        
        Args:
            slide_idx: 指定幻灯片索引（可选）
            include_notes: 是否包含备注
            
        Returns:
            文档的纯文本内容
        """
        if not self._is_loaded:
            self.load()
        
        slide_idx = kwargs.get('slide_idx')
        include_notes = kwargs.get('include_notes', False)
        
        texts = []
        
        slides_to_process = [slide_idx] if slide_idx is not None else range(self._slide_count)
        
        for idx in slides_to_process:
            if idx >= self._slide_count:
                continue
            
            slide = self._presentation.slides[idx]
            texts.append(f"=== Slide {idx + 1} ===")
            
            # 提取形状中的文本
            for shape in slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text.strip():
                    texts.append(shape_text)
            
            # 提取备注
            if include_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame.text.strip():
                    texts.append(f"[Notes]: {notes_text_frame.text}")
            
            texts.append("")  # 空行分隔
        
        return '\n'.join(texts)
    
    def _extract_shape_text(self, shape) -> str:
        """从形状中提取文本"""
        texts = []
        
        try:
            # 直接文本
            if hasattr(shape, 'text') and shape.text:
                texts.append(shape.text)
            
            # 文本框
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ' '.join(run.text for run in paragraph.runs)
                    if para_text.strip():
                        texts.append(para_text)
            
            # 表格
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_texts = []
                for row in shape.table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        table_texts.append(' | '.join(row_texts))
                if table_texts:
                    texts.append('\n'.join(table_texts))
            
            # 组合形状
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child_shape in shape.shapes:
                    child_text = self._extract_shape_text(child_shape)
                    if child_text.strip():
                        texts.append(child_text)
                        
        except Exception as e:
            logger.debug(f"提取形状文本失败: {e}")
        
        return '\n'.join(texts)
    
    def extract_content(self, **kwargs) -> ExtractedContent:
        """
        提取完整内容
        
        Args:
            include_images: 是否包含图片
            include_notes: 是否包含备注
            
        Returns:
            包含幻灯片内容的完整内容
        """
        if not self._is_loaded:
            self.load()
        
        include_images = kwargs.get('include_images', False)
        include_notes = kwargs.get('include_notes', False)
        
        content = ExtractedContent()
        
        for slide_idx, slide in enumerate(self._presentation.slides):
            slide_data = {
                'index': slide_idx,
                'slide_number': slide_idx + 1,
                'texts': [],
                'tables': [],
                'images': [],
                'notes': None
            }
            
            # 提取文本和表格
            for shape in slide.shapes:
                shape_text = self._extract_shape_text(shape)
                if shape_text.strip():
                    slide_data['texts'].append(shape_text)
                    content.paragraphs.append(shape_text)
                
                # 提取表格
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    if table_data:
                        slide_data['tables'].append(table_data)
                        content.tables.append(table_data)
                
                # 提取图片
                if include_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = self._extract_image_from_shape(shape, slide_idx)
                    if image_data:
                        slide_data['images'].append(image_data)
                        content.images.append(image_data)
            
            # 提取备注
            if include_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text
                if notes_text.strip():
                    slide_data['notes'] = notes_text
            
            content.slides.append(slide_data)
        
        content.text = '\n'.join(content.paragraphs)
        
        return content
    
    def _extract_image_from_shape(
        self,
        shape,
        slide_idx: int
    ) -> Optional[Dict[str, Any]]:
        """从形状中提取图片"""
        try:
            image = shape.image
            image_data = {
                'slide_index': slide_idx,
                'content_type': image.content_type,
                'ext': image.ext,
                'width': image.width,
                'height': image.height,
                'size': len(image.blob)
            }
            
            # 保存图片数据
            image_data['data'] = image.blob
            
            return image_data
            
        except Exception as e:
            logger.debug(f"提取图片失败: {e}")
            return None
    
    def extract_metadata(self, **kwargs) -> DocumentMetadata:
        """
        提取文档元数据
        
        Returns:
            文档元数据对象
        """
        if not self._is_loaded:
            self.load()
        
        metadata = DocumentMetadata()
        
        # 从文件系统获取基本元数据
        self._update_metadata_from_file()
        metadata.file_size = self._metadata.file_size
        metadata.file_path = self._metadata.file_path
        metadata.file_name = self._metadata.file_name
        metadata.file_extension = self._metadata.file_extension
        
        # 添加幻灯片信息
        metadata.custom_properties['slide_count'] = self._slide_count
        
        # 从core属性获取元数据
        if self._presentation:
            try:
                core_props = self._presentation.core_properties
                metadata.title = core_props.title
                metadata.author = core_props.author
                metadata.subject = core_props.subject
                metadata.keywords = core_props.keywords
                metadata.comments = core_props.comments
                metadata.last_modified_by = core_props.last_modified_by
                metadata.revision = core_props.revision
                metadata.category = core_props.category
                metadata.company = core_props.company
                
                # 日期处理
                if core_props.created:
                    metadata.created = core_props.created
                if core_props.modified:
                    metadata.modified = core_props.modified
                    
            except Exception as e:
                logger.debug(f"提取元数据失败: {e}")
        
        return metadata
    
    def save(self, output_path: Union[str, Path], **kwargs) -> "PowerPointDocument":
        """
        保存文档
        
        Args:
            output_path: 输出文件路径
            
        Returns:
            self: 支持链式调用
        """
        if not self._is_loaded:
            self.load()
        
        output_path = Path(output_path)
        
        if self._presentation:
            self._presentation.save(output_path)
        else:
            raise ValueError("没有可保存的演示文稿")
        
        return self
    
    def convert_to(
        self,
        target_type: DocumentType,
        output_path: Optional[Union[str, Path]] = None,
        **kwargs
    ) -> Union[str, Path, bytes]:
        """
        转换文档格式
        
        Args:
            target_type: 目标文档类型
            output_path: 输出路径（可选）
            
        Returns:
            转换后的文件路径或字节数据
        """
        if not self._is_loaded:
            self.load()
        
        if target_type == DocumentType.TXT:
            text = self.extract_text(**kwargs)
            if output_path:
                Path(output_path).write_text(text, encoding='utf-8')
                return output_path
            return text.encode('utf-8')
        
        elif target_type == DocumentType.PDF:
            return self._convert_to_pdf(output_path, **kwargs)
        
        elif target_type in [DocumentType.PPTX, DocumentType.PPT]:
            if output_path:
                self.save(output_path)
                return output_path
            else:
                buffer = io.BytesIO()
                self._presentation.save(buffer)
                return buffer.getvalue()
        
        else:
            raise ValueError(f"不支持的目标格式: {target_type}")
    
    def _convert_to_pdf(
        self,
        output_path: Optional[Union[str, Path]] = None,
        **kwargs
    ) -> Union[str, Path]:
        """
        转换为PDF
        
        Args:
            output_path: 输出路径
            
        Returns:
            PDF文件路径
        """
        if not output_path:
            output_path = self._temp_manager.create_temp_file(suffix='.pdf')
        else:
            output_path = Path(output_path)
        
        # 使用LibreOffice转换
        try:
            # 先保存为临时PPTX
            temp_pptx = self._temp_manager.create_temp_file(suffix='.pptx')
            self._presentation.save(temp_pptx)
            
            cmd = [
                'soffice',
                '--headless',
                '--convert-to', 'pdf',
                '--outdir', str(output_path.parent),
                str(temp_pptx)
            ]
            
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=120
            )
            
            if result.returncode != 0:
                raise RuntimeError(f"PDF转换失败: {result.stderr}")
            
            # 重命名输出文件
            expected_output = output_path.parent / f"{temp_pptx.stem}.pdf"
            if expected_output.exists() and expected_output != output_path:
                expected_output.rename(output_path)
            
            return output_path
            
        except Exception as e:
            logger.error(f"PDF转换失败: {e}")
            raise
    
    def extract_images(
        self,
        output_dir: Optional[Union[str, Path]] = None,
        **kwargs
    ) -> List[Path]:
        """
        提取所有图片
        
        Args:
            output_dir: 输出目录（可选）
            
        Returns:
            提取的图片路径列表
        """
        if not self._is_loaded:
            self.load()
        
        if not output_dir:
            output_dir = self._temp_manager.create_temp_dir()
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
        
        image_paths = []
        image_counter = 0
        
        for slide_idx, slide in enumerate(self._presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_counter += 1
                        
                        # 确定文件扩展名
                        ext = image.ext if image.ext else 'png'
                        image_path = output_dir / f"slide_{slide_idx + 1}_image_{image_counter}.{ext}"
                        
                        # 保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image.blob)
                        
                        image_paths.append(image_path)
                        
                    except Exception as e:
                        logger.warning(f"提取图片失败: {e}")
        
        return image_paths
    
    def close(self) -> None:
        """关闭文档，释放资源"""
        self._presentation = None
        self._temp_manager.cleanup()
        super().close()


class PowerPointHandler(BaseDocumentHandler):
    """
    PowerPoint文档处理器
    
    处理 .ppt 和 .pptx 格式的PowerPoint文档。
    """
    
    @property
    def supported_types(self) -> List[DocumentType]:
        """返回支持的文档类型列表"""
        return [DocumentType.PPT, DocumentType.PPTX]
    
    @property
    def supported_extensions(self) -> List[str]:
        """返回支持的文件扩展名列表"""
        return ['.ppt', '.pptx']
    
    def create_document(
        self,
        file_path: Optional[Union[str, Path]] = None,
        file_stream: Optional[BinaryIO] = None
    ) -> PowerPointDocument:
        """
        创建PowerPoint文档对象
        
        Args:
            file_path: 文件路径
            file_stream: 文件流
            
        Returns:
            PowerPoint文档对象
        """
        if file_path:
            self.validate_file(file_path)
        
        return PowerPointDocument(file_path=file_path, file_stream=file_stream)
    
    def create_new_document(self) -> PowerPointDocument:
        """
        创建新的空白PowerPoint文档
        
        Returns:
            新的PowerPoint文档对象
        """
        if not PPTX_AVAILABLE:
            raise ImportError("需要安装 python-pptx 库")
        
        doc = PowerPointDocument()
        doc._presentation = Presentation()
        doc._slide_count = 0
        doc._document_type = DocumentType.PPTX
        doc._is_loaded = True
        return doc


# 便捷函数
def extract_text_from_ppt(file_path: Union[str, Path], **kwargs) -> str:
    """
    从PowerPoint文档提取文本
    
    Args:
        file_path: PowerPoint文档路径
        **kwargs: 其他参数
        
    Returns:
        文档文本内容
    """
    handler = PowerPointHandler()
    doc = handler.create_document(file_path=file_path)
    with doc:
        return doc.extract_text(**kwargs)


def extract_images_from_ppt(
    file_path: Union[str, Path],
    output_dir: Optional[Union[str, Path]] = None,
    **kwargs
) -> List[Path]:
    """
    从PowerPoint文档提取图片
    
    Args:
        file_path: PowerPoint文档路径
        output_dir: 输出目录
        **kwargs: 其他参数
        
    Returns:
        图片路径列表
    """
    handler = PowerPointHandler()
    doc = handler.create_document(file_path=file_path)
    with doc:
        return doc.extract_images(output_dir=output_dir, **kwargs)
